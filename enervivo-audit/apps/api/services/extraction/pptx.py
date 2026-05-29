"""Extraction PPTX via python-pptx — texte des slides (titres, bullets, notes).

Cas d'usage typique : "Dossier de qualification J2a", "Présentation projet
mairie" — le contenu critique (titre, contexte, calendrier) est en texte
dans les slides. On extrait jusqu'à `_HEAD_BUDGET` chars puis on coupe.

Limite : si le PPTX est principalement composé d'images (screenshots de
plans), on récupère peu de texte → le LLM classera sur le filename + path,
ce qui reste utile.
"""

from __future__ import annotations

import io

from .base import HEAD_CHARS, ExtractionError, ScanNoTextError, TextExtractor

_HEAD_BUDGET = HEAD_CHARS * 2


class PptxExtractor(TextExtractor):
    def extract(self, content: bytes) -> str:
        try:
            from pptx import Presentation  # import paresseux
        except ImportError as e:  # pragma: no cover
            raise ExtractionError(f"python-pptx absent : {e}") from e

        try:
            prs = Presentation(io.BytesIO(content))
        except Exception as e:
            raise ExtractionError(f"python-pptx error : {e}") from e

        parts: list[str] = []
        total_len = 0
        for slide in prs.slides:
            slide_chunks: list[str] = []
            for shape in slide.shapes:
                # Bloc texte standard
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            slide_chunks.append(t)
                # Cellules de table
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                slide_chunks.append(t)
            # Notes
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_chunks.append(notes)
            except Exception:
                pass

            if slide_chunks:
                text = "\n".join(slide_chunks)
                parts.append(text)
                total_len += len(text)
                if total_len >= _HEAD_BUDGET:
                    break

        text = "\n\n---\n\n".join(parts)
        if not text.strip():
            # PPT sans texte (slides 100% image) — basculer sur vision LLM
            raise ScanNoTextError("PPTX sans texte (slides image uniquement)")
        return text
